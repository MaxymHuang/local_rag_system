"""
# File Finder RAG System
# Copyright (C) 2024 [Maxym (Yu-Chia) Huang]

# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:

# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.

# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.
"""

import os
from typing import List, Dict, Optional
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from pathlib import Path
import sys
import argparse
import time
import ollama
import json
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import requests

class FileSystemRAG:
    def __init__(self, root_dir: str = ".", ollama_host: str = "http://localhost:11434", 
                 ollama_model: str = "llama3.1:8b", sentence_model: str = "all-MiniLM-L6-v2"):
        self.root_dir = os.path.abspath(root_dir)  # Get absolute path
        # Check if path exists and is accessible
        if not os.path.exists(self.root_dir):
            raise ValueError(f"Directory does not exist: {self.root_dir}")
        if not os.access(self.root_dir, os.R_OK):
            raise ValueError(f"No read access to directory: {self.root_dir}")
            
        # Store the model name but don't load it yet (lazy loading)
        self.sentence_model_name = sentence_model
        self.model = None  # Will be loaded when needed
        self.index = None
        self.file_paths = []
        
        # Configure Ollama client
        self.ollama_host = ollama_host
        self.ollama_model = ollama_model
        
        # Test Ollama connection
        try:
            response = requests.get(f"{ollama_host}/api/tags")
            if response.status_code != 200:
                raise ConnectionError(f"Ollama server returned status code {response.status_code}")
            print("Successfully connected to Ollama server")
        except requests.exceptions.ConnectionError:
            raise ConnectionError(f"Could not connect to Ollama server at {ollama_host}. Make sure Ollama is running.")
        
    def _load_model_if_needed(self):
        """Load the sentence transformer model only when needed."""
        if self.model is None:
            print(f"Loading sentence transformer model: {self.sentence_model_name}")
            self.model = SentenceTransformer(self.sentence_model_name)
            print("Model loaded successfully!")
    
    def _get_file_description(self, file_path: str) -> str:
        """Generate a description for a file or directory."""
        path = Path(file_path)
        if path.is_dir():
            return f"Directory: {path.name} containing files and subdirectories"
        else:
            return f"File: {path.name} with extension {path.suffix}"
    
    def _read_file_contents(self, file_path: str) -> str:
        """Read the contents of a file."""
        path = Path(file_path)
        ext = path.suffix.lower()
        
        try:
            # Handle PDF files
            if ext == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf = PdfReader(file)
                    text = []
                    for page in pdf.pages:
                        text.append(page.extract_text())
                    return "\n".join(text)
            
            # Handle Word documents
            elif ext == '.docx':
                doc = Document(file_path)
                text = []
                for para in doc.paragraphs:
                    text.append(para.text)
                return "\n".join(text)
            
            # Handle PowerPoint files
            elif ext == '.pptx':
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        text.append("Slide: " + " | ".join(slide_text))
                return "\n".join(text)
            
            # Handle text files
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
                    
        except UnicodeDecodeError:
            return "Binary file - cannot be summarized"
        except Exception as e:
            return f"Error reading file: {str(e)}"
    
    def summarize_file(self, file_path: str, ollama_url: str = None, ollama_model: str = None) -> str:
        """Summarize a file using Ollama."""
        print(f"\nDebug: Attempting to summarize file: {file_path}")
        
        # Use provided settings or defaults
        ollama_url = ollama_url or self.ollama_host
        ollama_model = ollama_model or self.ollama_model
        
        if not os.path.isfile(file_path):
            print("Debug: Not a valid file")
            return "Not a file - cannot be summarized"
            
        # Check if Ollama server is available
        try:
            print(f"Debug: Checking Ollama server at {ollama_url}")
            response = requests.get(f"{ollama_url}/api/tags")
            if response.status_code != 200:
                print(f"Debug: Ollama server returned status code {response.status_code}")
                return f"Error: Ollama server returned status code {response.status_code}. Please ensure Ollama is running."
            print("Debug: Ollama server is available")
        except requests.exceptions.ConnectionError:
            print(f"Debug: Could not connect to Ollama server")
            return f"Error: Could not connect to Ollama server at {ollama_url}. Please ensure Ollama is running."
            
        content = self._read_file_contents(file_path)
        print(f"Debug: File content length: {len(content)} characters")
        if content.startswith("Error") or content.startswith("Binary"):
            print(f"Debug: File content indicates error or binary file")
            return content
            
        try:
            # Get file type for context
            file_type = Path(file_path).suffix.lower()
            file_name = Path(file_path).name
            print(f"Debug: Processing {file_type} file: {file_name}")
            
            # Prepare context-aware prompt
            if file_type == '.pdf':
                context = "PDF document"
            elif file_type == '.docx':
                context = "Word document"
            elif file_type == '.pptx':
                context = "PowerPoint presentation"
            else:
                context = "file"
                
            # Increase content length limit to 8000 characters
            prompt = f"""Please provide a concise summary of this {context} named '{file_name}'. 
IMPORTANT: Your response must be 600 words or less.

{content[:8000]}  # Increased content length limit

Focus on the main content and key points. Keep your summary under 600 words."""
            
            print(f"Debug: Using Ollama model: {ollama_model}")
            # Use Ollama to generate a summary
            try:
                print("Debug: Sending request to Ollama")
                response = ollama.chat(
                    model=ollama_model,
                    messages=[{
                        'role': 'user',
                        'content': prompt
                    }],
                    options={'host': ollama_url}
                )
                print("Debug: Received response from Ollama")
                print(f"Debug: Response type: {type(response)}")
                print(f"Debug: Response content: {response}")
                
                # Handle the Ollama response format - don't assume it's a dictionary
                summary = None
                
                # Try to extract message regardless of response type
                try:
                    # First try to get the message attribute/key
                    if hasattr(response, 'message'):
                        message = response.message
                        print(f"Debug: Found message attribute, type: {type(message)}")
                    elif isinstance(response, dict) and 'message' in response:
                        message = response['message']
                        print(f"Debug: Found message key in dict, type: {type(message)}")
                    else:
                        print(f"Debug: No message found in response")
                        print(f"Debug: Response attributes: {dir(response) if hasattr(response, '__dict__') else 'No attributes'}")
                        return "Error: No message found in Ollama response"
                    
                    print(f"Debug: Message type: {type(message)}")
                    print(f"Debug: Message content: {message}")
                    
                    # Now extract content from the message
                    if hasattr(message, 'content'):
                        summary = message.content
                        print(f"Debug: Successfully extracted content from Message object")
                    elif isinstance(message, dict) and 'content' in message:
                        summary = message['content']
                        print(f"Debug: Successfully extracted content from dictionary")
                    elif isinstance(message, str):
                        summary = message
                        print(f"Debug: Message is already a string")
                    else:
                        print(f"Debug: Message object has no 'content' attribute or key")
                        print(f"Debug: Message attributes: {dir(message) if hasattr(message, '__dict__') else 'No attributes'}")
                        return "Error: Message object missing content attribute"
                        
                except Exception as e:
                    print(f"Debug: Error extracting message: {str(e)}")
                    return f"Error: Failed to extract message from response: {str(e)}"
                
                if not summary:
                    print(f"Debug: Summary is empty or None")
                    return "Error: Empty response from Ollama server"
                
                print(f"Debug: Successfully extracted summary: {summary[:100]}...")
                
                # Count words and truncate if necessary
                words = summary.split()
                if len(words) > 600:
                    print(f"Debug: Truncating summary from {len(words)} to 600 words")
                    summary = ' '.join(words[:600]) + "..."
                
                return summary
                
            except requests.exceptions.ConnectionError:
                print("Debug: Lost connection to Ollama during chat")
                return f"Error: Lost connection to Ollama server. Please ensure Ollama is running."
            except Exception as e:
                print(f"Debug: Error during Ollama chat: {str(e)}")
                return f"Error generating summary: {str(e)}"
        except Exception as e:
            print(f"Debug: Error in summarize_file: {str(e)}")
            return f"Error generating summary: {str(e)}"
    
    def build_index(self):
        """Build the FAISS index from the file system."""
        # Load model only when building index
        self._load_model_if_needed()
        
        # Collect all files and directories
        self.file_paths = []
        start_time = time.time()
        files_processed = 0
        
        try:
            print("Scanning directory structure...")
            for root, dirs, files in os.walk(self.root_dir):
                # Skip hidden directories (starting with .)
                dirs[:] = [d for d in dirs if not d.startswith('.')]
                
                # Add directories
                for dir_name in dirs:
                    try:
                        full_path = os.path.normpath(os.path.join(root, dir_name))
                        self.file_paths.append(full_path)
                        files_processed += 1
                        if files_processed % 1000 == 0:
                            print(f"Processed {files_processed} items...")
                    except (PermissionError, OSError) as e:
                        print(f"Warning: Could not access {dir_name}: {str(e)}")
                        continue
                
                # Add files
                for file_name in files:
                    try:
                        if not file_name.startswith('.'):  # Skip hidden files
                            full_path = os.path.normpath(os.path.join(root, file_name))
                            self.file_paths.append(full_path)
                            files_processed += 1
                            if files_processed % 1000 == 0:
                                print(f"Processed {files_processed} items...")
                    except (PermissionError, OSError) as e:
                        print(f"Warning: Could not access {file_name}: {str(e)}")
                        continue
            
            if not self.file_paths:
                print("Warning: No files or directories found in the specified path.")
                return
            
            print(f"\nFound {len(self.file_paths)} files and directories.")
            print("Generating embeddings...")
            
            # Generate descriptions and embeddings
            descriptions = [self._get_file_description(path) for path in self.file_paths]
            embeddings = self.model.encode(descriptions)
            
            # Create FAISS index
            dimension = embeddings.shape[1]
            self.index = faiss.IndexFlatL2(dimension)
            self.index.add(embeddings.astype('float32'))
            
            end_time = time.time()
            print(f"\nIndex built successfully in {end_time - start_time:.2f} seconds!")
            
        except Exception as e:
            print(f"Error building index: {str(e)}")
            raise
    
    def search(self, query: str, k: int = 10) -> List[Dict[str, str]]:
        """Search for files/directories based on natural language query."""
        if self.index is None:
            raise ValueError("Index not built. Call build_index() first.")
        
        # Load model if not already loaded (should be loaded by now, but just in case)
        self._load_model_if_needed()
        
        # Encode query
        query_embedding = self.model.encode([query])
        
        # Search in FAISS index
        distances, indices = self.index.search(query_embedding.astype('float32'), k)
        
        # Return results
        results = []
        for i, idx in enumerate(indices[0]):
            if idx < len(self.file_paths):  # Ensure index is valid
                results.append({
                    'path': self.file_paths[idx],
                    'description': self._get_file_description(self.file_paths[idx]),
                    'relevance_score': float(1 / (1 + distances[0][i]))  # Convert distance to similarity score
                })
        
        return results

def main():
    # Set up argument parser
    parser = argparse.ArgumentParser(description='File Finder RAG System')
    parser.add_argument('--root-dir', type=str, default='.',
                      help='Root directory to search in (default: current directory)')
    parser.add_argument('--ollama-host', type=str, default='http://localhost:11434',
                      help='Ollama server host (default: http://localhost:11434)')
    parser.add_argument('--ollama-model', type=str, default='llama3.1:8b',
                      help='Ollama model name (default: llama3.1:8b)')
    parser.add_argument('--num-results', type=int, default=10,
                      help='Number of search results to return (default: 10)')
    args = parser.parse_args()

    try:
        # Initialize the RAG system with specified root directory and Ollama host
        print(f"Initializing file finder for directory: {os.path.abspath(args.root_dir)}")
        rag = FileSystemRAG(
            root_dir=args.root_dir, 
            ollama_host=args.ollama_host,
            ollama_model=args.ollama_model
        )
        
        # Build the index
        print("Building index...")
        rag.build_index()
        
        # Interactive search loop
        while True:
            try:
                query = input("\nEnter your search query (or 'quit' to exit): ")
                if query.lower() == 'quit':
                    break
                    
                results = rag.search(query, k=args.num_results)
                if not results:
                    print("\nNo results found.")
                    continue
                    
                print("\nSearch Results:")
                for i, result in enumerate(results, 1):
                    print(f"\n{i}. {result['description']}")
                    print(f"   Path: {result['path']}")
                    print(f"   Relevance Score: {result['relevance_score']:.2f}")
                
                # Ask if user wants to summarize a file
                summarize = input("\nWould you like to summarize any of these files? (y/n/q to quit): ").lower().strip()
                if summarize == 'q' or summarize == 'quit':
                    break
                if summarize != 'y':
                    continue
                
                # Ask user to select a file for summarization
                while True:
                    try:
                        selection = input("\nEnter the number of the file you want to summarize (or 'q' to quit, Enter to search again): ")
                        if not selection.strip():
                            break
                        if selection.lower() in ['q', 'quit']:
                            break
                            
                        idx = int(selection) - 1
                        if 0 <= idx < len(results):
                            selected_file = results[idx]['path']
                            print(f"\nGenerating summary for: {selected_file}")
                            summary = rag.summarize_file(selected_file)
                            print("\nSummary:")
                            print(summary)
                            break
                        else:
                            print("Invalid selection. Please enter a number from the list.")
                    except ValueError:
                        print("Please enter a valid number.")
                    except KeyboardInterrupt:
                        print("\nSelection cancelled.")
                        break
                        
            except KeyboardInterrupt:
                print("\nSearch interrupted. Type 'quit' to exit.")
                continue
            except Exception as e:
                print(f"\nError during search: {str(e)}")
                continue
                
    except Exception as e:
        print(f"Fatal error: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main() 